#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
遊戲家資訊科技 - 新人教育訓練簡報生成器
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_training_presentation():
    """創建新人教育訓練簡報"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 定義顏色
    primary_color = RGBColor(44, 62, 80)  # #2c3e50
    secondary_color = RGBColor(39, 172, 178)  # #27ACB2
    white = RGBColor(255, 255, 255)
    light_bg = RGBColor(248, 249, 250)
    
    # ========== 投影片 1：封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = light_bg
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "歡迎加入"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "遊戲家資訊科技有限公司"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(42)
    subtitle_para.font.bold = True
    subtitle_para.font.color.rgb = secondary_color
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 描述文字
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = "新人教育訓練"
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(32)
    desc_para.font.color.rgb = RGBColor(100, 100, 100)
    desc_para.alignment = PP_ALIGN.CENTER
    
    # ========== 投影片 2：公司概況 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    # 標題
    add_title(slide, "🏢 公司概況", primary_color)
    
    # 數據框（3個）
    stats = [
        ("105", "營運店家"),
        ("10萬", "會員數"),
        ("24/7", "無人營運")
    ]
    
    for i, (number, label) in enumerate(stats):
        left = Inches(1 + i * 2.8)
        top = Inches(1.8)
        width = Inches(2.3)
        height = Inches(1.5)
        
        # 背景框
        shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = 矩形
        shape.fill.solid()
        shape.fill.fore_color.rgb = secondary_color
        shape.line.fill.background()
        
        # 數字
        num_box = slide.shapes.add_textbox(left, top + Inches(0.3), width, Inches(0.6))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_para = num_frame.paragraphs[0]
        num_para.font.size = Pt(48)
        num_para.font.bold = True
        num_para.font.color.rgb = white
        num_para.alignment = PP_ALIGN.CENTER
        
        # 標籤
        label_box = slide.shapes.add_textbox(left, top + Inches(0.95), width, Inches(0.4))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(20)
        label_para.font.color.rgb = white
        label_para.alignment = PP_ALIGN.CENTER
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(3))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    add_bullet(content_frame, "ONE桌遊 - 自助桌遊館品牌", 22)
    add_bullet(content_frame, "全台最大自助麻將桌遊連鎖品牌", 18)
    add_bullet(content_frame, "AI 智能管理系統（門鎖、監控、繳費）", 18)
    add_bullet(content_frame, "24小時無人營運模式", 18)
    add_bullet(content_frame, "子品牌「智摸科技」- 麻將自助系統", 18)
    
    # ========== 投影片 3：品牌故事 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "📖 品牌故事", primary_color)
    
    timeline_data = [
        ("2023/7", "第一家店開幕（無人概念，尚未導入AI）\\n驗證市場需求與商業模式"),
        ("2023年底", "拓展至 6 家店\\n建立標準化營運流程"),
        ("2024", "快速擴張至 50 家店\\n開始研發智能管理系統"),
        ("2025", "突破 100 家店，累積 10萬會員\\n成立子品牌「智摸科技」"),
        ("2025/11", "正式導入 AI 技術\\n智能客服、選址系統、營運分析")
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for year, desc in timeline_data:
        p = content_frame.add_paragraph()
        p.text = f"{year}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = secondary_color
        p.space_before = Pt(12)
        
        p = content_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.space_after = Pt(8)
        p.level = 1
    
    # ========== 投影片 4：核心團隊 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "👥 核心團隊", primary_color)
    
    team_members = [
        ("👔", "建佑", "負責人"),
        ("💼", "Ken（阿建）", "總經理"),
        ("📱", "小圓", "業務"),
        ("💰", "小兔", "會計")
    ]
    
    for i, (icon, name, role) in enumerate(team_members):
        row = i // 2
        col = i % 2
        left = Inches(1.5 + col * 3.8)
        top = Inches(2 + row * 2.2)
        width = Inches(3.2)
        height = Inches(1.8)
        
        # 背景框
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = light_bg
        shape.line.color.rgb = secondary_color
        shape.line.width = Pt(2)
        
        # Icon
        icon_box = slide.shapes.add_textbox(left, top + Inches(0.2), width, Inches(0.5))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon
        icon_para = icon_frame.paragraphs[0]
        icon_para.font.size = Pt(40)
        icon_para.alignment = PP_ALIGN.CENTER
        
        # Name
        name_box = slide.shapes.add_textbox(left, top + Inches(0.8), width, Inches(0.4))
        name_frame = name_box.text_frame
        name_frame.text = name
        name_para = name_frame.paragraphs[0]
        name_para.font.size = Pt(22)
        name_para.font.bold = True
        name_para.font.color.rgb = primary_color
        name_para.alignment = PP_ALIGN.CENTER
        
        # Role
        role_box = slide.shapes.add_textbox(left, top + Inches(1.25), width, Inches(0.4))
        role_frame = role_box.text_frame
        role_frame.text = role
        role_para = role_frame.paragraphs[0]
        role_para.font.size = Pt(18)
        role_para.font.color.rgb = secondary_color
        role_para.alignment = PP_ALIGN.CENTER
    
    # ========== 投影片 5：AI 時代思維 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "🤖 AI 時代的轉型思維", primary_color)
    
    # 高亮框
    highlight_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.8), Inches(7), Inches(1.8))
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = secondary_color
    highlight_box.line.fill.background()
    
    highlight_text = slide.shapes.add_textbox(Inches(1.7), Inches(2), Inches(6.6), Inches(1.4))
    highlight_frame = highlight_text.text_frame
    highlight_frame.word_wrap = True
    
    p = highlight_frame.paragraphs[0]
    p.text = "當前時代背景"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = white
    
    add_bullet(highlight_frame, "我們正邁入 AI 時代（2-3 年轉型期）", 16, white)
    add_bullet(highlight_frame, "工具變化快，輔助能力強", 16, white)
    add_bullet(highlight_frame, "效率提升的速度遠超過去", 16, white)
    
    # 下方內容
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(7), Inches(2.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "必須習慣的改變"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    add_bullet(content_frame, "讓工具來協助工作，而非全靠人工", 18)
    add_bullet(content_frame, "從「我來做」轉變為「我來管理工具做」", 18)
    add_bullet(content_frame, "工作效率可以是過去的 10 倍甚至更多", 18)
    
    # ========== 投影片 6：駕馭 AI ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "🎯 駕馭 AI，而非被取代", primary_color)
    
    # 警告框
    warning_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.8), Inches(7), Inches(1.5))
    warning_box.fill.solid()
    warning_box.fill.fore_color.rgb = RGBColor(255, 243, 205)
    warning_box.line.color.rgb = RGBColor(255, 193, 7)
    warning_box.line.width = Pt(3)
    
    warning_text = slide.shapes.add_textbox(Inches(1.7), Inches(2), Inches(6.6), Inches(1.1))
    warning_frame = warning_text.text_frame
    warning_frame.word_wrap = True
    
    p = warning_frame.paragraphs[0]
    p.text = "⚠️ 關鍵觀念"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(133, 100, 4)
    
    p = warning_frame.add_paragraph()
    p.text = "不用擔心 AI 能取代掉誰"
    p.font.size = Pt(16)
    p.space_before = Pt(8)
    
    p = warning_frame.add_paragraph()
    p.text = "重點是如何駕馭 AI，創造更高效率及產值"
    p.font.size = Pt(16)
    p.font.bold = True
    
    # 高亮文字
    highlight_box2 = slide.shapes.add_shape(1, Inches(2), Inches(3.5), Inches(6), Inches(0.8))
    highlight_box2.fill.solid()
    highlight_box2.fill.fore_color.rgb = secondary_color
    highlight_box2.line.fill.background()
    
    highlight_text2 = slide.shapes.add_textbox(Inches(2.2), Inches(3.6), Inches(5.6), Inches(0.6))
    ht2_frame = highlight_text2.text_frame
    ht2_para = ht2_frame.paragraphs[0]
    ht2_para.text = "不改革，就會被改革淘汰"
    ht2_para.font.size = Pt(28)
    ht2_para.font.bold = True
    ht2_para.font.color.rgb = white
    ht2_para.alignment = PP_ALIGN.CENTER
    
    # 底部列表
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(7), Inches(2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    add_bullet(content_frame, "會用 AI 的人 > 不會用 AI 的人", 18)
    add_bullet(content_frame, "主動學習工具 > 固守舊方法", 18)
    add_bullet(content_frame, "擁抱變化 > 抗拒變化", 18)
    
    # ========== 投影片 7：公司發展階段 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "🚀 公司發展階段", primary_color)
    
    # 高亮框
    highlight_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.8), Inches(7), Inches(2.2))
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = secondary_color
    highlight_box.line.fill.background()
    
    highlight_text = slide.shapes.add_textbox(Inches(1.7), Inches(2), Inches(6.6), Inches(1.8))
    highlight_frame = highlight_text.text_frame
    highlight_frame.word_wrap = True
    
    p = highlight_frame.paragraphs[0]
    p.text = "我們正處於快速發展期"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = white
    
    add_bullet(highlight_frame, "從 1 家店 → 105 家店（不到 3 年）", 16, white)
    add_bullet(highlight_frame, "從無人概念 → AI 智能管理", 16, white)
    add_bullet(highlight_frame, "從單一品牌 → 多品牌布局（智摸科技）", 16, white)
    add_bullet(highlight_frame, "從手動管理 → 全自動化系統", 16, white)
    
    # 下方內容
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(7), Inches(2.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "我們需要什麼樣的夥伴？"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    add_bullet(content_frame, "願意共同衝刺的夥伴", 18)
    add_bullet(content_frame, "能快速適應變化的員工", 18)
    add_bullet(content_frame, "有主動學習精神的人", 18)
    add_bullet(content_frame, "能駕馭工具創造價值的人才", 18)
    
    # ========== 投影片 8：工作期待（重點） ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "💼 工作期待 - 醜話說在前面", primary_color)
    
    # 警告框
    warning_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.8), Inches(7), Inches(1.2))
    warning_box.fill.solid()
    warning_box.fill.fore_color.rgb = RGBColor(255, 243, 205)
    warning_box.line.color.rgb = RGBColor(255, 193, 7)
    warning_box.line.width = Pt(3)
    
    warning_text = slide.shapes.add_textbox(Inches(1.7), Inches(2), Inches(6.6), Inches(0.8))
    warning_frame = warning_text.text_frame
    
    p = warning_frame.paragraphs[0]
    p.text = "🚫 Ken（阿建）最不喜歡聽到的話："
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(133, 100, 4)
    
    p = warning_frame.add_paragraph()
    p.text = "「我不會」"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.alignment = PP_ALIGN.CENTER
    
    # 內容區塊
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.3), Inches(7), Inches(3.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "在 AI 時代，這些是基本能力"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    add_bullet(content_frame, "統整能力 - 整理資訊、歸納重點", 16)
    add_bullet(content_frame, "資料收集 - 搜尋、彙整、分析", 16)
    add_bullet(content_frame, "AI 工具可以做很完整的協助", 16)
    
    p = content_frame.add_paragraph()
    p.text = "我們需要的是什麼人？"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    p.space_before = Pt(15)
    
    add_bullet(content_frame, "✅ 會思考的人", 16)
    add_bullet(content_frame, "✅ 會整理的人", 16)
    add_bullet(content_frame, "✅ 會主動嘗試的人", 16)
    
    # ========== 投影片 9：工作態度 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "📋 正確的工作態度", primary_color)
    
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "收到任務後，你應該："
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    add_bullet(content_frame, "先簡單彙整任務（用工具協助）", 18)
    add_bullet(content_frame, "整理出初步結果（可以不完美）", 18)
    add_bullet(content_frame, "開會時討論與調整", 18)
    
    # 高亮框
    highlight_box = slide.shapes.add_shape(1, Inches(1.5), Inches(3.8), Inches(7), Inches(2.8))
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = secondary_color
    highlight_box.line.fill.background()
    
    highlight_text = slide.shapes.add_textbox(Inches(1.7), Inches(4), Inches(6.6), Inches(2.4))
    highlight_frame = highlight_text.text_frame
    highlight_frame.word_wrap = True
    
    p = highlight_frame.paragraphs[0]
    p.text = "可以接受的情況"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = white
    
    add_bullet(highlight_frame, "✅ 結果不盡理想 - 沒關係，可以調整", 16, white)
    add_bullet(highlight_frame, "✅ 判斷有些不同 - 沒關係，可以討論", 16, white)
    
    p = highlight_frame.add_paragraph()
    p.text = "無法接受的情況"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = white
    p.space_before = Pt(10)
    
    add_bullet(highlight_frame, "❌ 完全沒嘗試就說「我不會」", 16, white)
    add_bullet(highlight_frame, "❌ 不思考、不整理，等別人給答案", 16, white)
    
    # ========== 投影片 10：工作節奏 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "📅 工作節奏與溝通", primary_color)
    
    # 高亮框
    highlight_box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.8), Inches(7), Inches(2.2))
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = secondary_color
    highlight_box.line.fill.background()
    
    highlight_text = slide.shapes.add_textbox(Inches(1.7), Inches(2), Inches(6.6), Inches(1.8))
    highlight_frame = highlight_text.text_frame
    highlight_frame.word_wrap = True
    
    p = highlight_frame.paragraphs[0]
    p.text = "每週二下午 2:00 - 週會議"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    
    add_bullet(highlight_frame, "整理一週的工作狀況", 16, white)
    add_bullet(highlight_frame, "報告進度與問題", 16, white)
    add_bullet(highlight_frame, "討論解決方案", 16, white)
    add_bullet(highlight_frame, "設定下週目標", 16, white)
    
    # 下方內容
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(7), Inches(2.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "日常溝通方式"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    add_bullet(content_frame, "群組即時討論（快速決策）", 18)
    add_bullet(content_frame, "重要事項即時宣布", 18)
    add_bullet(content_frame, "⚠️ 重要事項開會需再重複一次", 18)
    add_bullet(content_frame, "確保所有人都吸收到訊息", 18)
    
    # ========== 投影片 11：常用工具 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "🔧 常用工具與系統", primary_color)
    
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "內部系統"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    add_bullet(content_frame, "ONE桌遊管理系統 - 店家管理、訂單、會員", 16)
    add_bullet(content_frame, "倉儲管理系統 - 庫存、進出貨", 16)
    add_bullet(content_frame, "AI 客服系統 - 自動回覆、問題分類", 16)
    add_bullet(content_frame, "打掃系統 - 任務派發、進度追蹤", 16)
    
    p = content_frame.add_paragraph()
    p.text = "AI 工具輔助"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    p.space_before = Pt(15)
    
    add_bullet(content_frame, "ChatGPT / Claude - 資料整理、文案撰寫", 16)
    add_bullet(content_frame, "Google Workspace - 文件協作", 16)
    add_bullet(content_frame, "LINE - 即時溝通", 16)
    add_bullet(content_frame, "GitHub - 程式碼管理（技術團隊）", 16)
    
    p = content_frame.add_paragraph()
    p.text = "💡 善用工具，讓效率翻倍！"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.space_before = Pt(15)
    
    # ========== 投影片 12：成長機會 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = white
    
    add_title(slide, "📈 職涯發展與成長機會", primary_color)
    
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "公司處於快速發展期，機會很多"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = secondary_color
    
    add_bullet(content_frame, "從執行 → 管理", 18)
    add_bullet(content_frame, "從單一職能 → 跨領域", 18)
    add_bullet(content_frame, "從傳統工作 → AI 協作", 18)
    add_bullet(content_frame, "從員工 → 核心夥伴", 18)
    
    # 高亮框
    highlight_box = slide.shapes.add_shape(1, Inches(1.5), Inches(4.2), Inches(7), Inches(2.2))
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = secondary_color
    highlight_box.line.fill.background()
    
    highlight_text = slide.shapes.add_textbox(Inches(1.7), Inches(4.4), Inches(6.6), Inches(1.8))
    highlight_frame = highlight_text.text_frame
    highlight_frame.word_wrap = True
    
    p = highlight_frame.paragraphs[0]
    p.text = "你將學到"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = white
    
    add_bullet(highlight_frame, "AI 時代的工作方法", 16, white)
    add_bullet(highlight_frame, "快速成長企業的營運模式", 16, white)
    add_bullet(highlight_frame, "多系統整合的實戰經驗", 16, white)
    add_bullet(highlight_frame, "從 0 到 1 建立新事業的能力", 16, white)
    
    # ========== 投影片 13：結語 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = light_bg
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "一起衝刺未來"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # 高亮框
    highlight_box = slide.shapes.add_shape(1, Inches(1.5), Inches(2.6), Inches(7), Inches(3))
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = secondary_color
    highlight_box.line.fill.background()
    
    highlight_text = slide.shapes.add_textbox(Inches(1.7), Inches(2.8), Inches(6.6), Inches(2.6))
    highlight_frame = highlight_text.text_frame
    highlight_frame.word_wrap = True
    
    p = highlight_frame.paragraphs[0]
    p.text = "你選擇加入我們，代表你認同我們的理念"
    p.font.size = Pt(18)
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    
    p = highlight_frame.add_paragraph()
    p.text = "這是一個快速成長的團隊"
    p.font.size = Pt(18)
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)
    
    p = highlight_frame.add_paragraph()
    p.text = "這是一個擁抱 AI 的時代"
    p.font.size = Pt(18)
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    
    p = highlight_frame.add_paragraph()
    p.text = "這是一個充滿機會的階段"
    p.font.size = Pt(18)
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    
    p = highlight_frame.add_paragraph()
    p.text = "讓我們一起駕馭工具、創造價值"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(15)
    
    # 底部文字
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.6))
    footer_frame = footer_box.text_frame
    footer_frame.text = "歡迎加入遊戲家！"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(32)
    footer_para.font.bold = True
    footer_para.font.color.rgb = secondary_color
    footer_para.alignment = PP_ALIGN.CENTER
    
    # 儲存簡報
    prs.save('遊戲家新人教育訓練.pptx')
    print("✅ PowerPoint 簡報已生成：遊戲家新人教育訓練.pptx")


def add_title(slide, text, color):
    """新增標題"""
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = color


def add_bullet(text_frame, text, font_size, color=None):
    """新增項目符號"""
    if color is None:
        color = RGBColor(80, 80, 80)
    
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.level = 1
    p.space_after = Pt(6)


if __name__ == '__main__':
    create_training_presentation()
