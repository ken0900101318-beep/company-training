#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
遊戲家資訊科技 - 新人教育訓練簡報生成器（完整版 v2.0）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """創建完整版簡報（25張投影片）"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 顏色定義
    primary = RGBColor(44, 62, 80)      # #2c3e50
    secondary = RGBColor(39, 172, 178)   # #27ACB2
    white = RGBColor(255, 255, 255)
    lightbg = RGBColor(248, 249, 250)
    warning_bg = RGBColor(255, 243, 205)
    warning_text = RGBColor(133, 100, 4)
    info_bg = RGBColor(209, 236, 241)
    
    # === 1. 封面 ===
    slide = add_blank_slide(prs, lightbg)
    add_text(slide, "歡迎加入", Inches(1), Inches(2), Inches(8), Inches(1), 
             size=60, bold=True, color=primary, align=PP_ALIGN.CENTER)
    add_text(slide, "遊戲家資訊科技有限公司", Inches(1), Inches(3.2), Inches(8), Inches(0.8),
             size=42, bold=True, color=secondary, align=PP_ALIGN.CENTER)
    add_text(slide, "新人教育訓練（完整版）", Inches(1), Inches(4.5), Inches(8), Inches(0.6),
             size=32, color=RGBColor(100,100,100), align=PP_ALIGN.CENTER)
    
    # === 2. 公司概況 ===
    slide = add_content_slide(prs, "🏢 公司概況")
    # 數據框
    stats = [("105", "營運店家"), ("10萬", "會員數"), ("24/7", "無人營運")]
    for i, (num, label) in enumerate(stats):
        left = Inches(1.5 + i * 2.5)
        add_stat_box(slide, left, Inches(1.8), Inches(2.2), Inches(1.5), num, label, secondary)
    
    add_bullets(slide, Inches(1.5), Inches(3.8), Inches(7), [
        "ONE桌遊 - 自助桌遊館品牌",
        "✅ 全台最大自助麻將桌遊連鎖品牌",
        "✅ AI 智能管理系統（門鎖、監控、繳費）",
        "✅ 24小時無人營運模式",
        "✅ 子品牌「智摸科技」- 麻將自助系統"
    ], 18)
    
    # === 3. 品牌故事 ===
    slide = add_content_slide(prs, "📖 品牌故事")
    timeline = [
        ("2023/7", "第一家店開幕（無人概念，尚未導入AI）\n驗證市場需求與商業模式"),
        ("2023年底", "拓展至 6 家店\n建立標準化營運流程"),
        ("2024", "快速擴張至 50 家店\n開始研發智能管理系統"),
        ("2025", "突破 100 家店，累積 10萬會員\n成立子品牌「智摸科技」"),
        ("2025/11", "正式導入 AI 技術\n智能客服、選址系統、營運分析")
    ]
    y = Inches(1.8)
    for year, content in timeline:
        add_text(slide, year, Inches(1.5), y, Inches(7), Inches(0.3), 
                 size=18, bold=True, color=secondary)
        add_text(slide, content, Inches(1.5), y + Inches(0.35), Inches(7), Inches(0.6),
                 size=14, color=RGBColor(80,80,80))
        y += Inches(1.1)
    
    # === 4. 核心團隊 ===
    slide = add_content_slide(prs, "👥 核心團隊")
    team = [
        ("👔", "建佑", "負責人"),
        ("💼", "Ken（阿建）", "總經理"),
        ("📱", "小圓", "業務"),
        ("💰", "小兔", "會計")
    ]
    for i, (icon, name, role) in enumerate(team):
        row, col = i // 2, i % 2
        left = Inches(1.5 + col * 3.8)
        top = Inches(2 + row * 2.2)
        add_team_card(slide, left, top, Inches(3.2), Inches(1.8), icon, name, role, lightbg, primary, secondary)
    
    # === 5. AI 時代思維 ===
    slide = add_content_slide(prs, "🤖 AI 時代的轉型思維")
    add_highlight_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(1.8), 
                      "當前時代背景", [
        "我們正邁入 AI 時代（2-3 年轉型期）",
        "工具變化快，輔助能力強",
        "效率提升的速度遠超過去"
    ], secondary)
    add_bullets(slide, Inches(1.5), Inches(4), Inches(7), [
        "必須習慣的改變",
        "🚀 讓工具來協助工作，而非全靠人工",
        "🚀 從「我來做」轉變為「我來管理工具做」",
        "🚀 工作效率可以是過去的 10 倍甚至更多"
    ], 18)
    
    # === 6. 駕馭 AI ===
    slide = add_content_slide(prs, "🎯 駕馭 AI，而非被取代")
    add_warning_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(1.5),
                    "⚠️ 關鍵觀念", [
        "不用擔心 AI 能取代掉誰",
        "重點是如何駕馭 AI，創造更高效率及產值"
    ])
    add_highlight_box(slide, Inches(2), Inches(3.5), Inches(6), Inches(0.8),
                      "不改革，就會被改革淘汰", [], secondary, center=True, large=True)
    add_bullets(slide, Inches(1.5), Inches(4.6), Inches(7), [
        "⭐ 會用 AI 的人 > 不會用 AI 的人",
        "⭐ 主動學習工具 > 固守舊方法",
        "⭐ 擁抱變化 > 抗拒變化"
    ], 18)
    
    # === 7. 公司發展階段 ===
    slide = add_content_slide(prs, "🚀 公司發展階段")
    add_highlight_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(2.2),
                      "我們正處於快速發展期", [
        "從 1 家店 → 105 家店（不到 3 年）",
        "從無人概念 → AI 智能管理",
        "從單一品牌 → 多品牌布局（智摸科技）",
        "從手動管理 → 全自動化系統"
    ], secondary)
    add_bullets(slide, Inches(1.5), Inches(4.3), Inches(7), [
        "我們需要什麼樣的夥伴？",
        "🚀 願意共同衝刺的夥伴",
        "🚀 能快速適應變化的員工",
        "🚀 有主動學習精神的人",
        "🚀 能駕馭工具創造價值的人才"
    ], 18)
    
    # === 8. 工作期待 ===
    slide = add_content_slide(prs, "💼 工作期待 - 醜話說在前面")
    add_warning_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(1.2),
                    "🚫 Ken（阿建）最不喜歡聽到的話", [
        "「我不會」"
    ], large_text=True)
    add_bullets(slide, Inches(1.5), Inches(3.3), Inches(7), [
        "在 AI 時代，這些是基本能力",
        "✅ 統整能力 - 整理資訊、歸納重點",
        "✅ 資料收集 - 搜尋、彙整、分析",
        "✅ AI 工具可以做很完整的協助"
    ], 16)
    add_highlight_box(slide, Inches(1.5), Inches(5), Inches(7), Inches(1.5),
                      "我們需要的是什麼人？", [
        "✅ 會思考的人",
        "✅ 會整理的人",
        "✅ 會主動嘗試的人"
    ], secondary)
    
    # === 9. 工作態度 ===
    slide = add_content_slide(prs, "📋 正確的工作態度")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "收到任務後，你應該：",
        "🎯 先簡單彙整任務（用工具協助）",
        "🎯 整理出初步結果（可以不完美）",
        "🎯 開會時討論與調整"
    ], 18)
    add_highlight_box(slide, Inches(1.5), Inches(3.8), Inches(7), Inches(2.8),
                      "可以接受 vs 無法接受", [
        "✅ 結果不盡理想 - 沒關係，可以調整",
        "✅ 判斷有些不同 - 沒關係，可以討論",
        "",
        "❌ 完全沒嘗試就說「我不會」",
        "❌ 不思考、不整理，等別人給答案"
    ], secondary)
    
    # === 10. 工作節奏 ===
    slide = add_content_slide(prs, "📅 工作節奏與溝通")
    add_highlight_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(2.2),
                      "每週二下午 2:00 - 週會議", [
        "📊 整理一週的工作狀況",
        "📢 報告進度與問題",
        "💬 討論解決方案",
        "🎯 設定下週目標"
    ], secondary, center_title=True)
    add_bullets(slide, Inches(1.5), Inches(4.3), Inches(7), [
        "日常溝通方式",
        "✅ 群組即時討論（快速決策）",
        "✅ 重要事項即時宣布",
        "⚠️ 重要事項開會需再重複一次",
        "🎯 確保所有人都吸收到訊息"
    ], 18)
    
    # === 11. 薪資福利制度 ===
    slide = add_content_slide(prs, "💰 薪資福利制度")
    add_info_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(0.8),
                 "💡 待填寫內容 - 請依實際情況填入具體數字與規定")
    add_bullets(slide, Inches(1.5), Inches(2.8), Inches(7), [
        "薪資結構",
        "💵 底薪：【待填寫】",
        "💵 績效獎金：【待填寫計算方式】",
        "💵 年終獎金：【待填寫規則】",
        "📅 發薪日：每月 __ 號",
        "",
        "保險福利",
        "✅ 勞保、健保（到職日起算）",
        "✅ 團體保險：【待填寫】",
        "",
        "休假制度",
        "🏖️ 特休：到職滿 __ 個月後",
        "🏥 病假、事假：【待填寫規定】",
        "📅 國定假日：依勞基法"
    ], 14)
    
    # === 12. 請假與加班 ===
    slide = add_content_slide(prs, "📝 請假與加班規定")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "請假流程",
        "1. 透過 【系統/LINE/紙本】 申請",
        "2. 需提前 __ 天（小時）",
        "3. 由 【直屬主管/Ken】 批准",
        "4. 緊急狀況請電話聯絡主管",
        "",
        "加班規定",
        "⏰ 加班費計算：【待填寫】",
        "🔄 補休規則：【待填寫】",
        "📋 加班需事前申請，主管同意後執行"
    ], 16)
    add_warning_box(slide, Inches(1.5), Inches(5.2), Inches(7), Inches(1),
                    "⚠️ 重要提醒", [
        "週會（每週二下午 2:00）為固定會議，若無法參加需提前告知"
    ])
    
    # === 13. 第一週工作計畫 ===
    slide = add_content_slide(prs, "📆 第一週工作計畫")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "Day 1 - 報到日",
        "• 09:00 報到、填寫資料",
        "• 10:00 辦公環境介紹",
        "• 11:00 系統帳號開通",
        "• 14:00 新人簡報講解（本簡報）",
        "• 16:00 認識團隊成員",
        "",
        "Day 2-3 - 熟悉階段",
        "• 系統操作教學",
        "• 跟著前輩觀察工作流程",
        "• 閱讀 SOP 文件",
        "• 提出疑問、記錄重點",
        "",
        "Day 4-5 - 實作階段",
        "• 開始簡單任務（有人帶著做）",
        "• 實際操作系統",
        "• 第一週回顧會議"
    ], 14)
    
    # === 14. 各職位職責與 KPI ===
    slide = add_content_slide(prs, "🎯 各職位職責與 KPI")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "業務（小圓）",
        "日常工作：",
        "  • 加盟諮詢接待",
        "  • 店家裝潢監工",
        "  • 材料訂購追蹤",
        "  • 客戶關係維護",
        "KPI 指標：",
        "  • 每月新簽約店數",
        "  • 裝潢準時完工率",
        "  • 客戶滿意度",
        "  • 【待補充】",
        "",
        "會計（小兔）",
        "日常工作：帳務處理、報表製作、發票開立、薪資計算",
        "KPI 指標：帳務正確率、報表準時繳交、【待補充】",
        "",
        "💡 具體 KPI 數字與考核標準將由主管說明"
    ], 13)
    
    # === 15. 系統操作教學 ===
    slide = add_content_slide(prs, "🔧 系統操作教學")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(3.5), [
        "內部系統",
        "🏪 ONE桌遊管理系統",
        "  店家管理、訂單、會員",
        "📦 倉儲管理系統",
        "  庫存、進出貨",
        "🤖 AI 客服系統",
        "  自動回覆、問題分類",
        "🧹 打掃系統",
        "  任務派發、進度追蹤"
    ], 14)
    add_bullets(slide, Inches(5.5), Inches(1.8), Inches(3.5), [
        "AI 工具輔助",
        "💬 ChatGPT / Claude",
        "  資料整理、文案撰寫",
        "📄 Google Workspace",
        "  文件協作",
        "📱 LINE",
        "  即時溝通",
        "💻 GitHub",
        "  程式碼管理（技術）"
    ], 14)
    add_info_box(slide, Inches(1.5), Inches(5.5), Inches(7), Inches(0.8),
                 "💡 帳號密碼由【IT/主管】提供，操作手冊請參考【Google Drive 連結】")
    
    # === 16. 緊急處理流程 ===
    slide = add_content_slide(prs, "🚨 緊急處理流程")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "店家緊急狀況",
        "• 門鎖故障（客戶鎖在裡面）→ 安撫客戶 → 遠端嘗試開鎖 → 通知業務",
        "• 設備故障（麻將桌、冷氣）→ 記錄問題 → 通知業務 → 派工維修",
        "• 客訴處理（鄰居噪音等）→ 了解情況 → 安撫客戶 → 通報主管",
        "• 系統當機 → 確認範圍 → 通知技術 → 暫用備案",
    ], 14)
    add_warning_box(slide, Inches(1.5), Inches(4.5), Inches(7), Inches(2),
                    "⚠️ 緊急聯絡人", [
        "Ken（總經理）- 【0912-XXX-XXX】",
        "建佑（負責人）- 【0912-XXX-XXX】",
        "小圓（業務）- 【0912-XXX-XXX】",
        "技術支援 - 【0912-XXX-XXX】"
    ])
    
    # === 17. 常見問題 FAQ ===
    slide = add_content_slide(prs, "❓ 常見問題 FAQ")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(3.5), [
        "關於工作",
        "Q: 忘記打卡怎麼辦？",
        "A: 【待填寫流程】",
        "",
        "Q: 臨時請假要找誰？",
        "A: 先告知直屬主管，",
        "   再填寫請假單",
        "",
        "Q: 系統登不進去找誰？",
        "A: 【IT 支援聯絡方式】",
        "",
        "Q: 週會沒空參加？",
        "A: 需提前告知 Ken，",
        "   並事後補看會議記錄"
    ], 12)
    add_bullets(slide, Inches(5.5), Inches(1.8), Inches(3.5), [
        "關於公司",
        "Q: 午餐公司有提供嗎？",
        "A: 【待填寫】",
        "",
        "Q: 停車位在哪？",
        "A: 【待填寫位置】",
        "",
        "Q: 可以在家工作嗎？",
        "A: 【待填寫遠端政策】",
        "",
        "Q: 多久會調薪？",
        "A: 【待填寫考核週期】"
    ], 12)
    
    # === 18. 成功/失敗案例 ===
    slide = add_content_slide(prs, "📚 成功/失敗案例分享")
    add_highlight_box(slide, Inches(1.5), Inches(1.8), Inches(7), Inches(1.2),
                      "✅ 成功案例：善用 AI 工具", [
        "情境：整理 100 家店的營收報表",
        "傳統做法：手動 Excel，需要 7 天",
        "AI 做法：用 Claude 寫腳本，3 天完成",
        "學到：善用工具可大幅提升效率"
    ], secondary, small_text=True)
    add_highlight_box(slide, Inches(1.5), Inches(3.2), Inches(7), Inches(1.2),
                      "✅ 成功案例：主動發現問題", [
        "情境：發現某供應商報價異常",
        "做法：主動比價、提出質疑",
        "結果：省下 50 萬成本",
        "學到：要有「老闆思維」"
    ], secondary, small_text=True)
    add_warning_box(slide, Inches(1.5), Inches(4.6), Inches(7), Inches(1.8),
                    "❌ 失敗案例", [
        "案例 1：說「我不會」導致延誤",
        "  情境：收到任務直接說不會，等主管教",
        "  結果：錯過截止日，影響整個專案",
        "  教訓：先嘗試用 AI/Google 找答案",
        "",
        "案例 2：溝通不確實導致誤會",
        "  教訓：重要事項要確認對方收到"
    ])
    
    # === 19. 公司願景與目標 ===
    slide = add_content_slide(prs, "🌟 公司願景與目標")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "短期目標（1 年內）",
        "📈 拓展到 150 家店",
        "🤖 AI 客服系統全面導入",
        "🎯 開發新產品線（除了麻將桌遊）",
        "",
        "中期目標（3 年內）",
        "🏆 成為全台最大自助娛樂品牌",
        "🌏 進軍海外市場（香港、新加坡）",
        "🎓 建立加盟主培訓學院"
    ], 16)
    add_highlight_box(slide, Inches(1.5), Inches(5.2), Inches(7), Inches(1.3),
                      "長期願景（5 年+）", [
        "✨ 重新定義休閒娛樂產業",
        "✨ AI 驅動的無人營運標準",
        "✨ 成為產業標竿"
    ], secondary)
    
    # === 20. 保密與資訊安全 ===
    slide = add_content_slide(prs, "🔒 保密與資訊安全")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "保密協議",
        "🔐 客戶資料不外流",
        "🔐 營收數據不公開",
        "🔐 AI 系統程式碼不外傳",
        "🔐 內部溝通記錄不截圖外傳",
        "",
        "資訊安全規範",
        "🔑 密碼設定：至少 12 字元，含英數符號",
        "🔑 禁止共用帳號",
        "🔑 離職時交還所有資料與權限",
        "🔑 LINE 群組截圖禁止外傳"
    ], 16)
    add_warning_box(slide, Inches(1.5), Inches(5.5), Inches(7), Inches(1),
                    "⚠️ 違規後果", [
        "第一次：書面警告 | 第二次：記過處分",
        "嚴重違規：立即解雇 + 法律追訴"
    ])
    
    # === 21. 辦公室實務資訊 ===
    slide = add_content_slide(prs, "🏢 辦公室實務資訊")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "上下班時間",
        "⏰ 正常工時：【9:00-18:00（待確認）】",
        "⏰ 彈性工時：【待填寫】",
        "🍱 午休時間：【12:00-13:00（待確認）】",
        "",
        "辦公地點",
        "📍 地址：【待填寫完整地址】",
        "🚗 停車資訊：【待填寫】",
        "🔑 門禁卡取得：【待填寫流程】",
        "",
        "辦公環境",
        "☕ 茶水間：【位置】",
        "🚻 洗手間：【位置】",
        "💻 會議室使用：【預約方式】",
        "♻️ 垃圾分類：【規則】"
    ], 15)
    
    # === 22. 職涯發展 ===
    slide = add_content_slide(prs, "📈 職涯發展與成長機會")
    add_bullets(slide, Inches(1.5), Inches(1.8), Inches(7), [
        "公司處於快速發展期，機會很多",
        "📊 從執行 → 管理",
        "📊 從單一職能 → 跨領域",
        "📊 從傳統工作 → AI 協作",
        "📊 從員工 → 核心夥伴"
    ], 18)
    add_highlight_box(slide, Inches(1.5), Inches(3.8), Inches(7), Inches(2.5),
                      "你將學到", [
        "🎯 AI 時代的工作方法",
        "🎯 快速成長企業的營運模式",
        "🎯 多系統整合的實戰經驗",
        "🎯 從 0 到 1 建立新事業的能力"
    ], secondary)
    
    # === 23. 公司文化 ===
    slide = add_content_slide(prs, "🌟 公司文化與價值觀")
    # 三個價值觀框
    values = [
        ("🚀", "快速行動", "想到就做\n邊做邊調整"),
        ("🤝", "團隊協作", "一起衝刺\n共享成果"),
        ("💡", "擁抱創新", "用 AI\n用新工具")
    ]
    for i, (icon, title, desc) in enumerate(values):
        left = Inches(1.5 + i * 2.5)
        add_value_box(slide, left, Inches(1.8), Inches(2.2), Inches(1.8), icon, title, desc, secondary)
    
    add_bullets(slide, Inches(1.5), Inches(4), Inches(7), [
        "我們相信",
        "⭐ 工具比經驗重要（AI 時代）",
        "⭐ 思考比執行重要（會思考才能駕馭工具）",
        "⭐ 速度比完美重要（快速迭代）",
        "⭐ 成長比穩定重要（發展期的選擇）"
    ], 18)
    
    # === 24. 結語 ===
    slide = add_blank_slide(prs, lightbg)
    add_text(slide, "一起衝刺未來", Inches(1), Inches(1.5), Inches(8), Inches(0.8),
             size=48, bold=True, color=primary, align=PP_ALIGN.CENTER)
    add_highlight_box(slide, Inches(1.5), Inches(2.6), Inches(7), Inches(3),
                      "", [
        "你選擇加入我們，代表你認同我們的理念",
        "",
        "這是一個快速成長的團隊",
        "這是一個擁抱 AI 的時代",
        "這是一個充滿機會的階段",
        "",
        "讓我們一起駕馭工具、創造價值"
    ], secondary, center_all=True)
    add_text(slide, "歡迎加入遊戲家！", Inches(1), Inches(6), Inches(8), Inches(0.6),
             size=32, bold=True, color=secondary, align=PP_ALIGN.CENTER)
    
    # === 25. 聯絡資訊 ===
    slide = add_content_slide(prs, "📞 聯絡資訊")
    add_bullets(slide, Inches(2), Inches(2), Inches(6), [
        "遊戲家資訊科技有限公司",
        "",
        "品牌：ONE桌遊 / 智摸科技",
        "營運店家：105 家",
        "會員數：10萬+",
        "",
        "地址：【待填寫】",
        "電話：【待填寫】",
        "Email：【待填寫】",
        "",
        "如有任何問題，歡迎隨時向直屬主管",
        "或 Ken（阿建）總經理詢問"
    ], 18)
    
    # 儲存
    prs.save('遊戲家新人教育訓練（完整版 v2.0）.pptx')
    print("✅ PowerPoint 完整版已生成（25張投影片）")

# === 輔助函數 ===

def add_blank_slide(prs, bg_color):
    """新增空白投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    return slide

def add_content_slide(prs, title):
    """新增內容投影片（白底 + 標題）"""
    slide = add_blank_slide(prs, RGBColor(255, 255, 255))
    add_text(slide, title, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
             size=32, bold=True, color=RGBColor(44, 62, 80))
    return slide

def add_text(slide, text, left, top, width, height, size=16, bold=False, color=None, align=PP_ALIGN.LEFT):
    """新增文字框"""
    if color is None:
        color = RGBColor(0, 0, 0)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb

def add_bullets(slide, left, top, width, items, size=16, color=None):
    """新增項目符號列表"""
    if color is None:
        color = RGBColor(80, 80, 80)
    tb = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0 if not item.startswith(('  ', '•', '  •')) else 1

def add_stat_box(slide, left, top, width, height, number, label, color):
    """新增數據框"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    
    add_text(slide, number, left, top + Inches(0.3), width, Inches(0.6),
             size=48, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + Inches(0.95), width, Inches(0.4),
             size=20, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)

def add_team_card(slide, left, top, width, height, icon, name, role, bg, name_color, role_color):
    """新增團隊成員卡片"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = role_color
    shape.line.width = Pt(2)
    
    add_text(slide, icon, left, top + Inches(0.2), width, Inches(0.5),
             size=40, align=PP_ALIGN.CENTER)
    add_text(slide, name, left, top + Inches(0.8), width, Inches(0.4),
             size=22, bold=True, color=name_color, align=PP_ALIGN.CENTER)
    add_text(slide, role, left, top + Inches(1.25), width, Inches(0.4),
             size=18, color=role_color, align=PP_ALIGN.CENTER)

def add_highlight_box(slide, left, top, width, height, title, items, color, center=False, large=False, center_title=False, small_text=False, center_all=False):
    """新增高亮框"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    
    y_offset = Inches(0.15)
    if title:
        title_size = 28 if large else 20
        title_align = PP_ALIGN.CENTER if center_title else PP_ALIGN.LEFT
        add_text(slide, title, left + Inches(0.2), top + y_offset, width - Inches(0.4), Inches(0.4),
                 size=title_size, bold=True, color=RGBColor(255,255,255), align=title_align)
        y_offset += Inches(0.5)
    
    if items:
        item_size = 14 if small_text else 16
        item_align = PP_ALIGN.CENTER if center_all else PP_ALIGN.LEFT
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + y_offset, width - Inches(0.4), height - y_offset - Inches(0.15))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(item_size)
            p.font.color.rgb = RGBColor(255,255,255)
            p.alignment = item_align

def add_warning_box(slide, left, top, width, height, title, items, large_text=False):
    """新增警告框"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 243, 205)
    shape.line.color.rgb = RGBColor(255, 193, 7)
    shape.line.width = Pt(3)
    
    add_text(slide, title, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
             size=20, bold=True, color=RGBColor(133, 100, 4))
    
    if items:
        item_size = 32 if large_text else 16
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(item_size)
            if large_text and i == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(192, 57, 43)
                p.alignment = PP_ALIGN.CENTER

def add_info_box(slide, left, top, width, height, text):
    """新增資訊框"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(209, 236, 241)
    shape.line.color.rgb = RGBColor(23, 162, 184)
    shape.line.width = Pt(3)
    
    add_text(slide, text, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3),
             size=16, color=RGBColor(12, 84, 96))

def add_value_box(slide, left, top, width, height, icon, title, desc, color):
    """新增價值觀框"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    
    add_text(slide, icon, left, top + Inches(0.15), width, Inches(0.5),
             size=40, align=PP_ALIGN.CENTER)
    add_text(slide, title, left, top + Inches(0.7), width, Inches(0.4),
             size=18, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
    add_text(slide, desc, left, top + Inches(1.15), width, Inches(0.5),
             size=14, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)

if __name__ == '__main__':
    create_presentation()
